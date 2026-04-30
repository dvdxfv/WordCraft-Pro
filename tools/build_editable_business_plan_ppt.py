from __future__ import annotations

import io
import re
import zipfile
from pathlib import Path
from typing import Iterable
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"G:\开发项目\wordcraft-pro")
OUT_DIR = ROOT / "artifacts" / "business-plan-editable"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPT = OUT_DIR / "business-plan-editable.pptx"

W = Inches(13.333)
H = Inches(7.5)


class C:
    BG = RGBColor(248, 250, 252)
    WHITE = RGBColor(255, 255, 255)
    NAVY = RGBColor(12, 38, 84)
    BLUE = RGBColor(37, 99, 235)
    TEAL = RGBColor(13, 148, 136)
    AMBER = RGBColor(245, 158, 11)
    SLATE = RGBColor(71, 85, 105)
    LIGHT = RGBColor(226, 232, 240)
    LIGHTER = RGBColor(239, 244, 248)
    SOFT_BLUE = RGBColor(219, 234, 254)
    SOFT_TEAL = RGBColor(204, 251, 241)
    SOFT_AMBER = RGBColor(254, 243, 199)


ANIMATIONS = {
    "appear": {"filter": None, "presetID": 1, "presetSubtype": 0},
    "fade": {"filter": "fade", "presetID": 10, "presetSubtype": 0},
    "wipe": {"filter": "wipe(left)", "presetID": 22, "presetSubtype": 1},
}


def add_textbox(slide, left, top, width, height, text, size=20, bold=False, color=C.NAVY,
               align=PP_ALIGN.LEFT, font_name="Microsoft YaHei", fill=None, margin=0.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin / 2)
    tf.margin_bottom = Inches(margin / 2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font_name
    f.color.rgb = color
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
    else:
        box.fill.background()
        box.line.fill.background()
    return box


def add_rect(slide, left, top, width, height, fill=C.WHITE, line=C.LIGHT, radius=True):
    shp_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.1)
    return shp


def add_placeholder(slide, left, top, width, height, title, subtitle):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.background()
    shp.line.color.rgb = C.BLUE
    shp.line.width = Pt(1.2)
    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(22)
    r1.font.bold = True
    r1.font.name = "Microsoft YaHei"
    r1.font.color.rgb = C.BLUE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(12)
    r2.font.name = "Microsoft YaHei"
    r2.font.color.rgb = C.SLATE
    return shp


def add_title(slide, title, subtitle):
    title_box = add_textbox(slide, Inches(0.6), Inches(0.28), Inches(6.8), Inches(0.7), title, size=34, bold=True)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.08), Inches(1.15), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C.BLUE
    accent.line.fill.background()
    subtitle_box = add_textbox(slide, Inches(0.62), Inches(1.16), Inches(7.3), Inches(0.42), subtitle, size=18, color=C.SLATE)
    return [title_box, accent, subtitle_box]


def add_footer_bar(slide, text, highlight=None):
    box = add_rect(slide, Inches(0.45), Inches(6.7), Inches(12.35), Inches(0.7), fill=C.WHITE, line=C.LIGHT)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    if highlight and highlight in text:
        before, after = text.split(highlight, 1)
        r.text = before
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.name = "Microsoft YaHei"
        r.font.color.rgb = C.NAVY
        r2 = p.add_run()
        r2.text = highlight
        r2.font.size = Pt(22)
        r2.font.bold = True
        r2.font.name = "Microsoft YaHei"
        r2.font.color.rgb = C.BLUE
        r3 = p.add_run()
        r3.text = after
        r3.font.size = Pt(22)
        r3.font.bold = True
        r3.font.name = "Microsoft YaHei"
        r3.font.color.rgb = C.NAVY
    else:
        r.text = text
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.name = "Microsoft YaHei"
        r.font.color.rgb = C.NAVY
    return box


def set_bg(slide, color=C.BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def card_with_text(slide, left, top, width, height, title, body, accent_color, idx=None):
    box = add_rect(slide, left, top, width, height)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    title_box = add_textbox(slide, left + Inches(0.28), top + Inches(0.14), width - Inches(0.55), Inches(0.35), title,
                            size=24, bold=True, color=accent_color)
    body_box = add_textbox(slide, left + Inches(0.28), top + Inches(0.52), width - Inches(0.55), Inches(0.42), body,
                           size=16, color=C.SLATE)
    shapes = [box, bar, title_box, body_box]
    if idx is not None:
        idx_box = add_textbox(slide, left + width - Inches(0.82), top + Inches(0.1), Inches(0.55), Inches(0.35), f"{idx:02d}",
                              size=20, bold=True, color=C.LIGHT, align=PP_ALIGN.RIGHT)
        shapes.append(idx_box)
    return shapes


def add_table_slide_table(slide, left, top, width, height, headers, rows, highlight_col=None):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    col_width = int(width / len(headers))
    for i in range(len(headers)):
        table.columns[i].width = col_width
    row_h = int(height / (len(rows) + 1))
    for i in range(len(rows) + 1):
        table.rows[i].height = row_h

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.BLUE if j == 0 else (C.TEAL if highlight_col == j else C.SOFT_BLUE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(18)
            r.font.name = "Microsoft YaHei"
            r.font.color.rgb = C.WHITE if (j == 0 or highlight_col == j) else C.NAVY

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            if highlight_col == j:
                cell.fill.fore_color.rgb = C.SOFT_TEAL
            else:
                cell.fill.fore_color.rgb = C.WHITE if i % 2 else C.LIGHTER
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(16)
                r.font.name = "Microsoft YaHei"
                r.font.bold = (j > 0)
                r.font.color.rgb = C.TEAL if highlight_col == j else (C.NAVY if j == 0 else C.SLATE)
    return slide.shapes[-1]


def create_sequence_timing_xml(targets: list[tuple[int, int, str]], duration: float = 0.35) -> str:
    if not targets:
        return ""
    dur_ms = int(duration * 1000)
    next_id = 3
    steps = []
    for shape_id, _delay_ms, animation in targets:
        info = ANIMATIONS.get(animation, ANIMATIONS["fade"])
        preset_id = info["presetID"]
        preset_subtype = info["presetSubtype"]
        wrapper_id, inner_id, leaf_id, set_id, eff_id = next_id, next_id + 1, next_id + 2, next_id + 3, next_id + 4
        next_id += 5
        if info["filter"] is None:
            effect_xml = (
                f'<p:set><p:cBhvr><p:cTn id="{set_id}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
                f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            )
        else:
            effect_xml = (
                f'<p:set><p:cBhvr><p:cTn id="{set_id}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
                f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
                f'<p:animEffect transition="in" filter="{escape(info["filter"])}"><p:cBhvr><p:cTn id="{eff_id}" dur="{dur_ms}"/>'
                f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl></p:cBhvr></p:animEffect>'
            )
        steps.append(
            f'<p:par><p:cTn id="{wrapper_id}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            f'<p:childTnLst><p:par><p:cTn id="{inner_id}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:par><p:cTn id="{leaf_id}" presetID="{preset_id}" presetClass="entr" presetSubtype="{preset_subtype}" fill="hold" nodeType="clickEffect">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>{effect_xml}</p:childTnLst>'
            f'</p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
        )
    bld_list = "".join(f'<p:bldP spid="{sid}" grpId="0"/>' for sid, _, _ in targets)
    return (
        '<p:timing><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        '<p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        + "".join(steps) +
        '</p:childTnLst></p:cTn><p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst><p:bldLst>'
        + bld_list +
        '</p:bldLst></p:timing>'
    )


def inject_animations(ppt_path: Path, slide_targets: dict[int, list[tuple[int, int, str]]]):
    raw = ppt_path.read_bytes()
    src = zipfile.ZipFile(io.BytesIO(raw), "r")
    out_mem = io.BytesIO()
    dst = zipfile.ZipFile(out_mem, "w", zipfile.ZIP_DEFLATED)
    for name in src.namelist():
        data = src.read(name)
        if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
            m = re.search(r"ppt/slides/slide(\d+)\.xml$", name)
            if not m:
                dst.writestr(name, data)
                continue
            slide_num = int(m.group(1))
            text = data.decode("utf-8")
            transition = '<p:transition advClick="1"><p:fade/></p:transition>'
            timing = create_sequence_timing_xml(slide_targets.get(slide_num, []))
            text = text.replace("</p:sld>", f"{transition}{timing}</p:sld>")
            data = text.encode("utf-8")
        dst.writestr(name, data)
    src.close()
    dst.close()
    ppt_path.write_bytes(out_mem.getvalue())


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    slide_targets: dict[int, list[tuple[int, int, str]]] = {}

    # Slide 1
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "WordCraft Pro 是什么", "不是泛 AI 写作工具，而是文档提交前检查产品")
    hero = add_rect(s, Inches(0.6), Inches(1.8), Inches(6.6), Inches(1.45), fill=C.WHITE, line=C.SOFT_BLUE)
    shapes.append(hero)
    shapes.append(add_textbox(s, Inches(0.9), Inches(2.05), Inches(6.0), Inches(0.95),
                              "帮助用户在提交前，\n提前发现最容易返工、退回、扣分的问题",
                              size=24, bold=True, color=C.BLUE))
    right = add_rect(s, Inches(8.1), Inches(1.35), Inches(4.45), Inches(4.35), fill=C.WHITE, line=C.LIGHT)
    shapes.append(right)
    shapes.append(add_textbox(s, Inches(8.5), Inches(1.7), Inches(3.7), Inches(0.45), "文档检查闭环", size=24, bold=True, color=C.TEAL, align=PP_ALIGN.CENTER))
    for i, t in enumerate(["文字校对", "格式检查", "交叉引用", "逻辑与表述优化"]):
        pill = add_rect(s, Inches(8.55), Inches(2.2 + i * 0.68), Inches(2.5), Inches(0.44), fill=C.SOFT_TEAL, line=C.SOFT_TEAL)
        shapes.append(pill)
        shapes.append(add_textbox(s, Inches(8.78), Inches(2.24 + i * 0.68), Inches(2.1), Inches(0.28), t, size=16, bold=True, color=C.TEAL))
    shapes.append(add_footer_bar(s, "对标 WPS 校对，但补上格式与引用检查", "格式与引用检查"))
    slide_targets[1] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 2
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "功能总览（上）", "先把 3 个核心能力讲清楚，再配产品截图")
    for idx, (title, body, color) in enumerate([
        ("上传预览", "支持 docx / pdf / xlsx，先打开文档再检查", C.BLUE),
        ("规则检查", "错别字、标点、一致性，先扫掉基础问题", C.TEAL),
        ("格式检查", "按学校、期刊、机构模板检查排版规范", C.AMBER),
    ], start=1):
        shapes += card_with_text(s, Inches(0.6), Inches(1.7 + (idx - 1) * 1.5), Inches(6.3), Inches(1.22), title, body, color, idx)
    for idx in range(3):
        ph = add_placeholder(s, Inches(7.45), Inches(1.55 + idx * 1.63), Inches(5.3), Inches(1.4), f"截图位 0{idx + 1}", "后续替换为产品截图")
        shapes.append(ph)
    shapes.append(add_footer_bar(s, "先让客户看到：我们不只是查文字", "我们不只是查文字"))
    slide_targets[2] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 3
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "功能总览（下）", "把最能拉开与 WPS 差距的能力放到这一页")
    shapes += card_with_text(s, Inches(0.6), Inches(1.9), Inches(5.8), Inches(1.5), "交叉引用",
                             "检查参考文献、图表编号、引用跳转是否断裂或错位", C.TEAL)
    shapes += card_with_text(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(1.5), "AI 深检",
                             "补上逻辑、语义、表述质量等高阶问题", C.BLUE)
    ph4 = add_placeholder(s, Inches(6.8), Inches(1.9), Inches(5.65), Inches(1.6), "截图位 04", "后续替换为产品截图")
    ph5 = add_placeholder(s, Inches(6.8), Inches(4.0), Inches(5.65), Inches(1.6), "截图位 05", "后续替换为产品截图")
    shapes += [ph4, ph5]
    shapes.append(add_footer_bar(s, "WPS 更偏文字校对，我们补上提交前最关键的格式与引用检查", "我们补上"))
    slide_targets[3] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 4
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "为什么对标 WPS", "WPS 教育了用户“校对有价值”，但还没有覆盖最痛的提交前问题")
    headers = ["能力项", "WPS 校对", "WordCraft Pro"]
    rows = [
        ["错别字 / 标点 / 语病", "强项，偏文字层", "能做，而且能串联整份文档检查"],
        ["排版格式检查", "模板化能力弱", "核心创新点之一"],
        ["交叉引用检查", "基本不解决", "核心创新点之一"],
        ["使用场景", "写作过程中的文字校对", "正式提交前的文档验收"],
    ]
    tbl = add_table_slide_table(s, Inches(0.65), Inches(1.75), Inches(12.0), Inches(4.55), headers, rows, highlight_col=2)
    shapes.append(tbl)
    shapes.append(add_footer_bar(s, "我们不是替代 WPS，而是补上它没有解决的高成本问题", "补上它没有解决的高成本问题"))
    slide_targets[4] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 5
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "我们的核心创新点", "不是“我们也有 AI”，而是“我们能查出别人查不出的提交前问题”")
    cards = [
        ("排版格式检查", "学校、期刊、机构最在意的规范问题", C.BLUE),
        ("交叉引用检查", "参考文献、图表编号、跳转错位", C.TEAL),
        ("规则层 + AI 层组合", "确定性问题更稳，理解性问题更深", C.AMBER),
        ("模板与规则沉淀", "从个人复用，到团队共享，再到机构规则库", C.BLUE),
    ]
    positions = [(0.65, 1.9), (6.65, 1.9), (0.65, 4.0), (6.65, 4.0)]
    for (title, body, color), (x, y) in zip(cards, positions):
        shapes += card_with_text(s, Inches(x), Inches(y), Inches(5.35), Inches(1.35), title, body, color)
    shapes.append(add_footer_bar(s, "创新点不是功能堆叠，而是更贴近提交前的真实痛点"))
    slide_targets[5] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 6
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "谁会先买单", "先从最刚需、最容易成交的人群切入，再往上走")
    audiences = [
        ("学生 / 研究生", "离提交最近，返工痛感最强\n最适合做第一波增长", C.BLUE),
        ("研究团队 / 咨询团队", "多人协作后，最怕交付不统一\n最适合做 Team 版本切入", C.TEAL),
        ("高校 / 机构", "看重规范统一和规则沉淀\n客单价更高，但成交节奏更慢", C.AMBER),
    ]
    for i, (title, body, color) in enumerate(audiences):
        x = 0.65 + i * 4.1
        shapes += card_with_text(s, Inches(x), Inches(2.0), Inches(3.75), Inches(2.6), title, body, color)
    shapes.append(add_footer_bar(s, "先起量，再提客单价，再做机构化"))
    slide_targets[6] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 7
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "怎么收费", "先让客户用起来，再让客户为深度和效率付费")
    plans = [
        ("Free", ["先体验基础检查", "少量 AI 次数", "先看到价值"], C.BLUE, C.SOFT_BLUE),
        ("Pro", ["AI 深检", "更多额度", "为个人效率付费"], C.TEAL, C.SOFT_TEAL),
        ("Team", ["规则共享", "批量检查", "为协作交付付费"], C.AMBER, C.SOFT_AMBER),
    ]
    for i, (name, bullets, color, bg) in enumerate(plans):
        x = 0.85 + i * 4.05
        card = add_rect(s, Inches(x), Inches(1.9), Inches(3.5), Inches(3.6), fill=C.WHITE, line=color)
        shapes.append(card)
        shapes.append(add_textbox(s, Inches(x + 0.3), Inches(2.25), Inches(2.9), Inches(0.55), name, size=30, bold=True, color=color, align=PP_ALIGN.CENTER))
        for j, bullet in enumerate(bullets):
            shapes.append(add_textbox(s, Inches(x + 0.35), Inches(3.0 + j * 0.52), Inches(2.9), Inches(0.34), f"• {bullet}", size=18, color=C.SLATE))
    shapes.append(add_textbox(s, Inches(1.9), Inches(5.85), Inches(9.8), Inches(0.35), "体验  →  深度  →  协作", size=20, bold=True, color=C.TEAL, align=PP_ALIGN.CENTER))
    shapes.append(add_footer_bar(s, "不是先卖会员，而是先让客户看到问题和价值"))
    slide_targets[7] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 8
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "未来 12 个月目标", "用三档情景对齐注册、付费、收入和留存目标")
    headers = ["指标", "保守", "基准", "乐观"]
    rows = [
        ["注册用户", "300", "600", "1200"],
        ["付费用户（月底）", "60", "150", "350"],
        ["MRR（第12月）", "¥594", "¥1485", "¥3465"],
        ["ARR 折算", "¥7128", "¥17820", "¥41580"],
        ["月均留存率（付费）", "55%", "65%", "75%"],
        ["团队版客户", "0", "2", "5"],
        ["定制服务收入（全年）", "¥3000", "¥12000", "¥30000"],
    ]
    tbl = add_table_slide_table(s, Inches(0.48), Inches(1.72), Inches(12.15), Inches(4.65), headers, rows, highlight_col=2)
    shapes.append(tbl)
    shapes.append(add_footer_bar(s, "基准情景是默认执行线，保守线用于止损，乐观线用于评估上限"))
    slide_targets[8] = [(sh.shape_id, 0, "fade") for sh in shapes]

    # Slide 9
    s = prs.slides.add_slide(blank)
    set_bg(s)
    shapes = []
    shapes += add_title(s, "建议价区间与首发价", "长期价锚定价值，首发价服务冷启动转化")
    headers = ["套餐", "建议长期价", "首发价（冷启动期）", "首发期结束条件"]
    rows = [
        ["Pro 月付", "¥9.9/月", "¥6.9/首月", "付费用户满50人"],
        ["Pro 年付", "¥79/年", "¥69/年", "付费用户满50人"],
        ["Team 月付", "¥49/月", "¥39/月", "团队客户满3个"],
        ["Team 年付", "¥399/年", "¥299/年", "团队客户满3个"],
    ]
    tbl = add_table_slide_table(s, Inches(0.48), Inches(1.72), Inches(12.15), Inches(4.55), headers, rows, highlight_col=2)
    shapes.append(tbl)
    shapes.append(add_footer_bar(s, "首发价用于验证付费闭环，不是长期永久低价"))
    slide_targets[9] = [(sh.shape_id, 0, "fade") for sh in shapes]

    prs.save(OUT_PPT)
    inject_animations(OUT_PPT, slide_targets)
    print(OUT_PPT)


if __name__ == "__main__":
    build()
